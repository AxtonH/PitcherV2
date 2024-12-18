import os

from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()


import streamlit as st
import google.generativeai as genai

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import io
import re
from reportlab.pdfgen import canvas  # To generate PDFs
from reportlab.lib.pagesizes import letter

# Google Gemini API Configuration
GEMINI_API_KEY = st.secrets["API_KEY"]
genai.configure(api_key=GEMINI_API_KEY)

# Constants for Safe Content Length
SAFE_CHAR_LIMIT = 15000  # Approximate safe character limit for Gemini API input

# Function to summarize content
def summarize_text(content):
    try:
        model = genai.GenerativeModel("gemini-pro")
        response = model.generate_content(f"Summarize this text briefly:\n\n{content}")
        return response.text.strip()
    except Exception as e:
        return f"Error summarizing text: {e}"

# Function to process text chunks
def process_text_chunks(chunks):
    total_length = sum(len(chunk) for chunk in chunks)
    if total_length > SAFE_CHAR_LIMIT:
        return [summarize_text("\n".join(chunks))]
    return chunks

# Function to extract text from files
def extract_text_from_file(file):
    if file.type == "text/plain":
        return file.read().decode("utf-8")
    elif file.type == "application/pdf":
        pdf_reader = PdfReader(file)
        return "\n".join([page.extract_text() or "" for page in pdf_reader.pages])
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = Document(file)
        return "\n".join([para.text for para in doc.paragraphs])
    return None

# Function to chunk text into logical segments
def chunk_text(text, max_length=500):
    sentences = re.split(r'(?<=[.!?]) +', text)
    chunks = []
    current_chunk = ""
    for sentence in sentences:
        if len(current_chunk) + len(sentence) <= max_length:
            current_chunk += f"{sentence} "
        else:
            chunks.append(current_chunk.strip())
            current_chunk = sentence
    if current_chunk:
        chunks.append(current_chunk.strip())
    return chunks

# Define templates for Pitch Deck and Corporate Portfolio
def get_template(context_type):
    if context_type == "Pitch Deck":
        return ["Introduction", "Problem Statement", "Solution", "Key Features", "Market Opportunity", "Financials"]
    elif context_type == "Corporate Portfolio":
        return ["Company Overview", "Services/Products", "Achievements", "Team", "Contact Details"]
    return []

# Improved function to parse AI response into slides
def parse_response_to_slides(ai_response, template):
    slides = []
    slide_number = 1

    # Remove asterisks and leading dashes from AI response
    cleaned_response = re.sub(r'\*\*|\*', '', ai_response)  # Remove asterisks
    cleaned_response = re.sub(r'^- *', '', cleaned_response, flags=re.MULTILINE)  # Remove dashes

    # Create a regex pattern to detect template titles
    title_patterns = "|".join([re.escape(title) for title in template])
    section_splits = re.split(f"({title_patterns})", cleaned_response)

    # Iterate through template titles and match content dynamically
    for idx, title in enumerate(template):
        content = ""
        # Look for the title in the split sections
        for i in range(1, len(section_splits), 2):  # Check every second item (titles)
            if section_splits[i].strip().lower() == title.lower():
                content = section_splits[i + 1].strip() if i + 1 < len(section_splits) else ""
                break

        # Process and clean up content into bullet points
        content_lines = []
        for line in content.split("\n"):
            line = line.strip()
            if line:
                content_lines.append(f"• {line}")

        # Add slide with title and cleaned content
        slides.append({
            "title": f"{slide_number}. {title}",
            "content": "\n".join(content_lines) or "• No content available."
        })
        slide_number += 1

    return slides

# Function to export slides into PowerPoint
def export_slides(slides):
    presentation = Presentation()
    for slide_info in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_info.get("title", "Slide Title")
        content_shape = slide.placeholders[1]
        content_shape.text = slide_info.get("content", "")
    return presentation

# Function to export slides into TXT
def export_slides_to_txt(slides):
    txt_content = ""
    for slide in slides:
        txt_content += f"{slide['title']}\n{slide['content']}\n\n"
    return txt_content

# Function to export slides into PDF
def export_slides_to_pdf(slides):
    pdf_bytes = io.BytesIO()
    c = canvas.Canvas(pdf_bytes, pagesize=letter)
    width, height = letter

    for slide in slides:
        # Write slide title
        c.setFont("Helvetica-Bold", 14)
        c.drawString(72, height - 50, slide['title'])

        # Write slide content
        c.setFont("Helvetica", 11)
        content_lines = slide['content'].split("\n")
        y = height - 80
        for line in content_lines:
            c.drawString(72, y, line)
            y -= 15  # Move to the next line
            if y < 50:  # Add new page if out of space
                c.showPage()
                y = height - 50

        c.showPage()  # End current slide's page

    c.save()
    pdf_bytes.seek(0)
    return pdf_bytes

# Streamlit App
st.set_page_config(page_title="🚀 Deck Generator & Chatbot", layout="centered")
st.title("🚀 Professional Presentation Generator & Chatbot")
st.caption("Turn your documents into engaging presentations and chat with the assistant.")

# Initialize session state variables
if "slides" not in st.session_state:
    st.session_state.slides = []
if "current_presentation" not in st.session_state:
    st.session_state.current_presentation = None
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

# Slide Editor Functionality
def display_slide_editor():
    if st.session_state.slides:
        st.subheader("📋 Review and Edit Slides")
        for idx, slide in enumerate(st.session_state.slides):
            with st.expander(f"{slide['title']}"):
                new_title = st.text_input(f"Slide {idx + 1} Title", slide["title"], key=f"title_{idx}")
                new_content = st.text_area(f"Slide {idx + 1} Content", slide["content"], height=150,
                                           key=f"content_{idx}")

                if st.button(f"Save Changes to Slide {idx + 1}", key=f"save_{idx}"):
                    slide["title"] = new_title
                    slide["content"] = new_content
                    st.session_state.current_presentation = export_slides(st.session_state.slides)
                    st.success("Slide updated successfully!")

    if st.session_state.current_presentation:
        col1, col2, col3 = st.columns(3)  # Create 3 columns

        with col1:  # Download PPTX button in the first column
            ppt_bytes = io.BytesIO()
            st.session_state.current_presentation.save(ppt_bytes)
            ppt_bytes.seek(0)
            st.download_button(
                label="💾 Download Presentation as PPT",
                data=ppt_bytes,
                file_name="generated_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        with col2:  # Download PDF button in the second column
            pdf_bytes = export_slides_to_pdf(st.session_state.slides)
            st.download_button(
                label="📄 Download Presentation as PDF",
                data=pdf_bytes,
                file_name="generated_presentation.pdf",
                mime="application/pdf"
            )

        with col3:  # Download TXT button in the third column
            txt_content = export_slides_to_txt(st.session_state.slides)
            txt_bytes = io.BytesIO(txt_content.encode("utf-8"))
            st.download_button(
                label="📝 Download Presentation as TXT",
                data=txt_bytes,
                file_name="generated_presentation.txt",
                mime="text/plain"
            )

# Radio button for context selection
context_type = st.radio(
    "What would you like to do?",
    options=["Pitch Deck", "Corporate Portfolio", "Chatbot"],
    index=0,
    horizontal=True
)

# Context Logic
if context_type != "Chatbot":
    st.subheader("📂 Upload Document")
    uploaded_file = st.file_uploader("Choose a file (PDF, Word, or TXT)", type=["txt", "pdf", "docx"])

    if uploaded_file:
        file_content = extract_text_from_file(uploaded_file)
        if file_content:
            text_chunks = chunk_text(file_content)
            summarized_chunks = process_text_chunks(text_chunks)

            if st.button("Generate Presentation"):
                with st.spinner("Generating slides..."):
                    system_context = (
                        f"You are an AI assistant tasked with generating a professional {context_type}. "
                        f"Create detailed content for the following sections: {', '.join(get_template(context_type))}. "
                        "For each section, provide a clear heading followed by 3-5 bullet points."
                    )
                    model = genai.GenerativeModel("gemini-pro")
                    response = model.generate_content(
                        f"{system_context}\n\nHere is the input text:\n{'\n'.join(summarized_chunks)}"
                    )

                    slides = parse_response_to_slides(response.text, get_template(context_type))
                    st.session_state.slides = slides
                    st.session_state.current_presentation = export_slides(slides)
                    st.success(f"{context_type} slides generated successfully!")

    display_slide_editor()

else:
    st.subheader("💬 Chat with Assistant")
    user_input = st.chat_input("Type your message:")
    if user_input:
        st.session_state.chat_history.append({"role": "user", "content": user_input})

        with st.spinner("Assistant is typing..."):
            try:
                model = genai.GenerativeModel("gemini-pro")
                history = "\n".join(
                    [f"{msg['role'].capitalize()}: {msg['content']}" for msg in st.session_state.chat_history])
                response = model.generate_content(history)
                reply = response.text.strip()
                st.session_state.chat_history.append({"role": "assistant", "content": reply})
            except Exception as e:
                reply = f"Error: {e}"
                st.session_state.chat_history.append({"role": "assistant", "content": reply})

    if st.session_state.chat_history:
        for chat in st.session_state.chat_history:
            with st.chat_message(chat["role"]):
                st.markdown(chat["content"])
